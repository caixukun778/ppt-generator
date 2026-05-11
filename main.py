import os
import uuid
import json
import shutil
import time
import httpx
from fastapi import FastAPI, HTTPException, Request, UploadFile, File as FastAPIFile
from fastapi.responses import FileResponse
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

app = FastAPI()
OUTPUT_DIR = "outputs"
os.makedirs(OUTPUT_DIR, exist_ok=True)

# ========== 颜色映射 ==========
COLOR_MAP = {
    "深蓝色": RGBColor(0, 51, 102),
    "蓝色": RGBColor(0, 102, 204),
    "科技蓝": RGBColor(0, 153, 255),
    "灰色": RGBColor(128, 128, 128),
    "深灰色": RGBColor(64, 64, 64),
    "白色": RGBColor(255, 255, 255),
    "黑色": RGBColor(0, 0, 0),
}

def get_color(name: str) -> RGBColor:
    for key, value in COLOR_MAP.items():
        if key in name:
            return value
    return RGBColor(0, 51, 102)


# ========== PPT 生成核心 ==========
def build_pptx(title: str, style_desc: str, slides_json_str: str):
    """生成 PPT，并内建专业排版规则"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1. 解析风格
    primary_color = get_color(style_desc)
    accent_color = get_color(style_desc)  # 可扩展

    # 2. 解析内容
    slides_data = []
    try:
        # 清洗可能的 markdown 标记
        cleaned = slides_json_str.strip()
        if cleaned.startswith("```json"):
            cleaned = cleaned[7:]
        if cleaned.endswith("```"):
            cleaned = cleaned[:-3]
        slides_data = json.loads(cleaned)
        if not isinstance(slides_data, list):
            slides_data = [{"title": "内容", "type": "content", "blocks": []}]
    except:
        slides_data = [{"title": title, "type": "content", "blocks": [{"heading": "内容生成中", "bullets": ["请稍后重试"]}]}]

    # 3. 逐页生成，应用排版规则
    for slide_data in slides_data:
        slide_type = slide_data.get("type", "content")

        if slide_type == "cover":
            # ---------- 封面排版 ----------
            slide_layout = prs.slide_layouts[6]  # 空白版式
            slide = prs.slides.add_slide(slide_layout)

            # 全屏背景色
            bg = slide.background
            bg.fill.solid()
            bg.fill.fore_color.rgb = primary_color

            # 大标题（居中）
            txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10), Inches(2))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = slide_data.get("title", title)
            p.font.size = Pt(48)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER

            # 副标题（下对齐）
            sub = slide_data.get("subtitle", "")
            if sub:
                txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(5), Inches(10), Inches(1.2))
                tf2 = txBox2.text_frame
                p2 = tf2.paragraphs[0]
                p2.text = sub
                p2.font.size = Pt(24)
                p2.font.color.rgb = RGBColor(220, 220, 220)
                p2.alignment = PP_ALIGN.CENTER

        elif slide_type == "toc":
            # ---------- 目录排版 ----------
            slide_layout = prs.slide_layouts[1]  # 标题+内容
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = slide_data.get("title", "目录")

            body = slide.shapes.placeholders[1].text_frame
            body.clear()
            items = slide_data.get("items", [])
            for i, item in enumerate(items):
                p = body.add_paragraph()
                # 自动识别图表/案例标签，增加视觉区分
                if "📊" in item or "💡" in item:
                    p.text = f"    {item}"
                    p.font.color.rgb = accent_color
                    p.font.size = Pt(18)
                else:
                    p.text = f"{i+1}.  {item}"
                    p.font.size = Pt(20)
                p.space_after = Pt(10)

        else:
            # ---------- 内容页排版（左右分栏） ----------
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = slide_data.get("title", "内容")

            # 左栏：文字要点
            left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.5), Inches(5))
            left_tf = left_box.text_frame
            left_tf.word_wrap = True

            # 右栏：图表/案例卡片
            right_box = slide.shapes.add_textbox(Inches(7.5), Inches(1.8), Inches(5.5), Inches(5))
            right_tf = right_box.text_frame
            right_tf.word_wrap = True

            blocks = slide_data.get("blocks", [])
            has_right_content = False

            for block in blocks:
                heading = block.get("heading", "")
                bullets = block.get("bullets", [])

                # 如果是图表建议或案例，放进右栏
                if "📊" in heading or "💡" in heading:
                    p = right_tf.add_paragraph()
                    p.text = heading
                    p.font.size = Pt(18)
                    p.font.bold = True
                    p.font.color.rgb = accent_color
                    for bullet in bullets:
                        p2 = right_tf.add_paragraph()
                        p2.text = f"• {bullet}"
                        p2.font.size = Pt(14)
                        p2.font.color.rgb = RGBColor(80, 80, 80)
                        p2.space_after = Pt(6)
                    has_right_content = True
                else:
                    # 普通文字块放进左栏
                    p = left_tf.add_paragraph()
                    p.text = heading
                    p.font.size = Pt(20)
                    p.font.bold = True
                    p.font.color.rgb = primary_color
                    for bullet in bullets:
                        p2 = left_tf.add_paragraph()
                        p2.text = f"• {bullet}"
                        p2.font.size = Pt(16)
                        p2.space_after = Pt(4)

            # 如果没有右栏内容，给右栏一个提示
            if not has_right_content:
                p = right_tf.add_paragraph()
                p.text = "📌 建议在此补充图表或案例"
                p.font.size = Pt(16)
                p.font.italic = True
                p.font.color.rgb = RGBColor(150, 150, 150)

    # 4. 保存文件
    filename = f"{uuid.uuid4()}.pptx"
    filepath = os.path.join(OUTPUT_DIR, filename)
    prs.save(filepath)
    return filepath, filename

@app.post("/generate_pptx")
async def generate_pptx(request: Request):
    try:
        data = await request.json()
        title = data.get("title", "PPT")
        style_desc = data.get("style", "")
        slides_str = data.get("slides", "[]")
        filepath, filename = build_pptx(title, style_desc, slides_str)
        download_url = f"/download/{filename}"
        return {"success": True, "download_url": download_url}
    except Exception as e:
        return {"success": False, "error": str(e)}


@app.get("/download/{filename}")
async def download_file(filename: str):
    filepath = os.path.join(OUTPUT_DIR, filename)
    if not os.path.exists(filepath):
        raise HTTPException(status_code=404, detail="File not found or expired")
    return FileResponse(filepath, filename=filename)


# ========== 素材库 ==========
UPLOAD_DIR = "uploads"
os.makedirs(UPLOAD_DIR, exist_ok=True)

IMAGES_DB = "images_db.json"
if not os.path.exists(IMAGES_DB):
    with open(IMAGES_DB, "w", encoding="utf-8") as f:
        json.dump([], f)

@app.post("/upload_image")
async def upload_image(request: Request):
    """上传图片：支持直接传文件 + 传图片链接，两者可同时进行"""
    saved = []
    tag = ""

    try:
        content_type = request.headers.get("content-type", "")

        # 1. 处理文件上传（multipart/form-data）
        if "multipart" in content_type:
            form = await request.form()
            tag = form.get("tag", "")
            file = form.get("file")
            if file and hasattr(file, "filename"):
                ext = file.filename.split(".")[-1] if "." in file.filename else "png"
                save_name = f"{uuid.uuid4()}.{ext}"
                save_path = os.path.join(UPLOAD_DIR, save_name)
                with open(save_path, "wb") as f:
                    f.write(await file.read())
                saved.append({"name": save_name, "url": f"/uploads/{save_name}"})

        # 2. 处理 JSON 链接上传（可以单独，也可以和文件流同时存在）
        # 如果content-type是application/json，或者form里包含image_url
        if "json" in content_type:
            data = await request.json()
            tag = data.get("tag", tag)
            image_url = data.get("image_url", "")
            if image_url:
                async with httpx.AsyncClient() as client:
                    resp = await client.get(image_url)
                    if resp.status_code == 200:
                        ct = resp.headers.get("content-type", "")
                        ext = "png"
                        if "jpeg" in ct or "jpg" in ct:
                            ext = "jpg"
                        elif "webp" in ct:
                            ext = "webp"
                        save_name = f"{uuid.uuid4()}.{ext}"
                        save_path = os.path.join(UPLOAD_DIR, save_name)
                        with open(save_path, "wb") as f:
                            f.write(resp.content)
                        saved.append({"name": save_name, "url": f"/uploads/{save_name}"})
        else:
            # 如果multipart里也传了image_url（非标准但可能发生）
            form = await request.form()
            tag = tag or form.get("tag", "")
            image_url = form.get("image_url", "")
            if image_url:
                async with httpx.AsyncClient() as client:
                    resp = await client.get(str(image_url))
                    if resp.status_code == 200:
                        ct = resp.headers.get("content-type", "")
                        ext = "png"
                        if "jpeg" in ct or "jpg" in ct:
                            ext = "jpg"
                        elif "webp" in ct:
                            ext = "webp"
                        save_name = f"{uuid.uuid4()}.{ext}"
                        save_path = os.path.join(UPLOAD_DIR, save_name)
                        with open(save_path, "wb") as f:
                            f.write(resp.content)
                        saved.append({"name": save_name, "url": f"/uploads/{save_name}"})

        if not saved:
            return {"success": False, "error": "未提供图片文件或链接"}

        # 3. 保存记录
        with open(IMAGES_DB, "r", encoding="utf-8") as f:
            images = json.load(f)

        for s in saved:
            images.append({
                "id": s["name"],
                "tag": tag,
                "url": s["url"],
                "uploaded_at": time.strftime("%Y-%m-%d %H:%M:%S")
            })

        with open(IMAGES_DB, "w", encoding="utf-8") as f:
            json.dump(images, f, ensure_ascii=False, indent=2)

        return {
            "success": True,
            "saved": len(saved),
            "images": [{"id": s["name"], "url": s["url"]} for s in saved]
        }

    except Exception as e:
        return {"success": False, "error": str(e)}

@app.get("/list_images")
async def list_images(tag: str = ""):
    """查看图片列表，可按标签筛选"""
    try:
        with open(IMAGES_DB, "r", encoding="utf-8") as f:
            images = json.load(f)

        if tag:
            images = [img for img in images if img.get("tag") == tag]

        for img in images:
            img["full_url"] = f"https://ppt-generator-production-a9fd.up.railway.app{img['url']}"

        return {"success": True, "images": images, "count": len(images)}
    except Exception as e:
        return {"success": False, "error": str(e)}


@app.get("/uploads/{filename}")
async def get_uploaded_file(filename: str):
    """访问上传的图片"""
    filepath = os.path.join(UPLOAD_DIR, filename)
    if not os.path.exists(filepath):
        raise HTTPException(status_code=404, detail="Image not found")
    return FileResponse(filepath)


@app.get("/cleanup")
def cleanup():
    """重置图片数据库"""
    try:
        with open(IMAGES_DB, "w", encoding="utf-8") as f:
            json.dump([], f)
        return {"success": True}
    except Exception as e:
        return {"success": False, "error": str(e)}


@app.get("/")
def root():
    return {"status": "running"}
